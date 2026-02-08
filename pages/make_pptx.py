import sys
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt(md_path, ppt_path):
    if not os.path.exists(md_path):
        print(f"Error: {md_path} not found.")
        return

    prs = Presentation()
    
    # Define a simple "Dark Theme"
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(13, 17, 23) # Dark background matching app

    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by Slide headers ### [Slide X]
    slides_content = re.split(r'### \[Slide \d+\]', content)
    
    # First part is the report overview
    intro_match = re.search(r'# \[(.*?)\](.*?)\n(.*?)\n', slides_content[0], re.S)
    if intro_match:
        title_text = intro_match.group(2).strip()
        subtitle_text = intro_match.group(3).strip()
        
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55) # Gold
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    # Process each slide
    for slide_md in slides_content[1:]:
        lines = slide_md.strip().split('\n')
        if not lines: continue
        
        slide_title = lines[0].strip()
        body_lines = [l.strip() for l in lines[1:] if l.strip()]
        
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_dark_background(slide)
        
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55)
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for line in body_lines:
            p = tf.add_paragraph()
            clean_line = line.replace('**', '').replace('- ', '').replace('• ', '').replace('✅ ', 'V ')
            p.text = clean_line
            p.font.color.rgb = RGBColor(230, 230, 230)
            p.font.size = Pt(18)
            
            if line.startswith('- ') or line.startswith('• ') or line.startswith('✅ '):
                p.level = 1
            else:
                p.level = 0
                p.font.bold = True

    prs.save(ppt_path)
    print(f"PPT saved to {ppt_path}")

if __name__ == "__main__":
    if len(sys.argv) > 2:
        md_file = sys.argv[1]
        ppt_file = sys.argv[2]
    else:
        # Default fallback
        md_file = r"c:\Users\PEACE\5차_AI\Lotte_AI_Browser_RunBAT_Demo\LotteTower_AI_SalesApp_Python\수익_극대화_전략_리포트_v4_0.md"
        ppt_file = r"c:\Users\PEACE\5차_AI\Lotte_AI_Browser_RunBAT_Demo\LotteTower_AI_SalesApp_Python\pages\수익_극대화_전략_리포트_v4_0.pptx"
    
    create_ppt(md_file, ppt_file)
